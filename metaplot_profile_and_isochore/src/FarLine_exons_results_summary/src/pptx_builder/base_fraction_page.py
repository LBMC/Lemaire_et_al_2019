#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from math import floor

## base_fraction page
inter_img_space_y = Inches( 0.1 )
offset_x = Inches( 0.25 )
base_fraction_img_height = floor( 2 * img_height / 3 )
second_column_x = prs.slide_width - slide_marge_right - 2 * base_fraction_img_height

# reg_dic = { 1: [ '50', '100' ], 2: [ '200', '500' ] } #~2: [ '50', '200' ] }
reg_dic = { 1: [ in_exon_list[ 0 ], out_off_list[ 0 ] ], 2: [ in_exon_list[ 1 ], out_off_list[ 1 ] ] } #~2: [ '50', '200' ] }
for index in reg_dic.keys():
    ## initialize position of the cursor
    current_y = slide_marge_top
    current_x = slide_marge_left

    page_title = 'base_fractions'
    file_path = script_dir + '/slide_creator.py'
    exec( exec_text )

    ## exon_scheme
    inEx = reg_dic[ index ][ 0 ] #'50'
    outOff = reg_dic[ index ][ 1 ] #'100'
    file_path = script_dir + '/splicing_site_region_scheme.py'
    exec( exec_text )


    base_fraction_dir = fig_dir + '/base_fraction/inEx%s_outOff%s/win20b' % ( reg_dic[ index ][ 0 ], reg_dic[ index ][ 1 ] )

    fig_part_y = current_y + height_tbx + Inches( 0.5 )

    ######
    # simple base fraction
    #   # splicing sites titles
    current_y = fig_part_y
    textbox = slide.shapes.add_textbox( left=current_x + offset_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '3\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    textbox = slide.shapes.add_textbox( left=current_x + offset_x + base_fraction_img_height, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '5\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   # base A
    current_y = current_y + height_tbx
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( A )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_A.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_A.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base T
    current_y = current_y + img.height + inter_img_space_y
    current_x = slide_marge_left
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( T )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_T.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_T.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base C
    current_y = current_y + img.height + inter_img_space_y
    current_x = slide_marge_left
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( C )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_C.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_C.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base G
    current_y = current_y + img.height + inter_img_space_y
    current_x = slide_marge_left
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( G )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_G.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_G.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    ######
    # double base fraction
    #   # splicing sites titles
    current_x = second_column_x
    current_y = fig_part_y
    textbox = slide.shapes.add_textbox( left=current_x + offset_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '3\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    textbox = slide.shapes.add_textbox( left=current_x + offset_x + base_fraction_img_height, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '5\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   # base CG
    current_y = current_y + height_tbx
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( GC )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_GC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_GC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base AT
    current_y = current_y + img.height + inter_img_space_y
    current_x = second_column_x
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( AT )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_AT.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_AT.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base TC
    current_y = current_y + img.height + inter_img_space_y
    current_x = second_column_x
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( TC )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_TC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_TC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )


    #   # base AC
    current_y = current_y + img.height + inter_img_space_y
    current_x = second_column_x
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '( AC )'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = base_fraction_dir + '/' + exon_list_name + '_3SS_AC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    #   #   # 5'SS
    current_x = current_x + img.width
    image_file = base_fraction_dir + '/' + exon_list_name + '_5SS_AC.png'
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=base_fraction_img_height )

    pass
