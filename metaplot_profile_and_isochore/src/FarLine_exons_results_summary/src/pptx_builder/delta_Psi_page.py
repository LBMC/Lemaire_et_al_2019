#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from math import floor

## delta_Psi page
inter_img_space_y = Inches( 0.1 )
inter_img_space_x = Inches( 0.1 )
offset_x = Inches( 0.25 )
delta_Psi_img_height = floor( 2 * img_height / 3 )
inter_line_space_y = delta_Psi_img_height
# second_column_x = prs.slide_width - slide_marge_right - 2 * delta_Psi_img_height

delta_Psi_dir = fig_dir + '/deltaPSI_plot'

fig_part_y = slide_marge_top + height_tbx + Inches( 0.5 )
first_col_x = slide_marge_left + height_tbx + Inches( 1 )
second_col_x = first_col_x + offset_x + delta_Psi_img_height
third_col_x = second_col_x + offset_x + delta_Psi_img_height

## load table of figure names
with open( delta_Psi_dir + '/fig_name_table.tsv') as fich:
    fig_name_table = fich.readlines()
    pass

fig_name_list = [''] * len( fig_name_table )
for idx,line in enumerate( fig_name_table ):
    fig_name_list[ idx ] = line[0:-1].split('\t')
    pass

######
# simple base fraction
page_nbr = 0
line_nbr = 0
for item in fig_name_list:
    if line_nbr % 5 == 0:
        page_nbr += 1

        ## initialize position of the cursor
        current_y = slide_marge_top
        current_x = slide_marge_left

        ## create the slide
        page_title = 'delta_Psi ( page %d )' % page_nbr
        file_path = script_dir + '/slide_creator.py'
        exec( exec_text )

        #   # deltaPsi plot type titles
        current_y = fig_part_y
        current_x = first_col_x
        textbox = slide.shapes.add_textbox( left=current_x + offset_x, top=current_y, width=Inches( 0.95 ), height=height_tbx )
        textbox.text = 'all exons'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

        current_x = second_col_x
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.6 ), height=height_tbx )
        textbox.text = 'correlated exons'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

        current_x = third_col_x
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
        textbox.text = 'anti-correlated exons'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

        current_y = current_y + height_tbx - inter_line_space_y
        pass

    line_nbr += 1

    nickname = item[ 0 ]
    fig_names = item[ 1 ].split(',')

    #   # line of figures
    current_y = current_y + inter_line_space_y
    current_x = slide_marge_left
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.3 ), height=height_tbx )
    textbox.text = nickname
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # all exons
    current_y = current_y + height_tbx
    current_x = first_col_x
    image_file = fig_names[ 0 ]
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=delta_Psi_img_height )

    #   #   # correlated exons
    current_x = second_col_x
    image_file = fig_names[ 1 ]
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=delta_Psi_img_height )

    #   #   # anti-correlated exons
    current_x = third_col_x
    image_file = fig_names[ 2 ]
    img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=delta_Psi_img_height )

    pass
