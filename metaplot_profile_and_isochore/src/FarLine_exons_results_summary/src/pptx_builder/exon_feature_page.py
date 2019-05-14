#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN

## exon_features page
## initialize position of the cursor
current_y = slide_marge_top
current_x = slide_marge_left

page_title = 'exon_features'
file_path = script_dir + '/slide_creator.py'
exec( exec_text )

## exon_scheme
file_path = script_dir + '/exon_scheme.py'
exec( exec_text )


exon_feature_dir = fig_dir + '/exon_feature'
FarLine_psi_dir = fig_dir + '/FarLine_psi'

# exon length plot
current_y = current_y + height_tbx + Inches( 1 )
textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
textbox.text = 'exon_length'
textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

current_y = current_y + height_tbx
image_file = exon_feature_dir + '/' + exon_list_name + '_exon_length_boxplot.png'
img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=img_height )


# intron length plot
current_y = current_y - height_tbx
current_x = current_x + img.width + Inches( 0.5 )
textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
textbox.text = 'intron_length'
textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

current_y = current_y + height_tbx
image_file = exon_feature_dir + '/' + exon_list_name + '_intron_length_boxplot.png'
img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=img_height )


# acceptor strength
current_y = current_y + img_height + Inches( 0.5 )
current_x = slide_marge_left
textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
textbox.text = '3\'SS strength'
textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

current_y = current_y + height_tbx
image_file = exon_feature_dir + '/' + exon_list_name + '_force_acceptor_boxplot.png'
img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=img_height )


# donor strength
current_y = current_y - height_tbx
current_x = current_x + img.width + Inches( 0.5 )
textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
textbox.text = '5\'SS strength'
textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

current_y = current_y + height_tbx
image_file = exon_feature_dir + '/' + exon_list_name + '_force_donor_boxplot.png'
img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=img_height )


# siPP/siGL2 PSI boxplot
current_y = current_y + img_height + Inches( 0.5 )
current_x = slide_marge_left
textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 2 ), height=height_tbx )
textbox.text = 'siPP and siGL2 PSI'
textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

current_y = current_y + height_tbx
image_file = FarLine_psi_dir + '/' + exon_list_name + '_siPP-siGL2_PSI_boxplot.png'
img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=img_height )

# # counts of exons
# nbr_psi_100_95 = int( getoutput( 'wc -l %s/psi_100_95_FDB_list.tsv' % exon_list_dir ).split()[0] ) - 1
# nbr_psi_60_40 = int( getoutput( 'wc -l %s/psi_60_40_FDB_list.tsv' % exon_list_dir ).split()[0] ) - 1
# nbr_psi_5_0 = int( getoutput( 'wc -l %s/psi_5_0_FDB_list.tsv' % exon_list_dir ).split()[0] ) - 1
# nbr_up_exons = int( getoutput( 'wc -l %s/exon_up_%s_list.tsv' % ( exon_list_dir, exon_list_name ) ).split()[0] ) - 1
# nbr_down_exons = int( getoutput( 'wc -l %s/exon_down_%s_list.tsv' % ( exon_list_dir, exon_list_name ) ).split()[0] ) - 1
#
# current_x += img.width + Inches( 1.5 )
# textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 3 ), height=height_tbx )
# textbox.text = '%d\tpsi_100_95\n%d\tpsi_60_40\n%d\tpsi_5_0\n%d\tup exons\n%d\tdown exons' % ( nbr_psi_100_95, nbr_psi_60_40, nbr_psi_5_0, nbr_up_exons, nbr_down_exons )
# textbox.text_frame.paragraphs[0].font.size = font_size_big
