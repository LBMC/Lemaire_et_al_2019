#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
import sys
from os import makedirs,path

image_dir = sys.argv[ 1 ] #~ '../kmeansp2_clustering_refSignal'
#~_img_base = heatmap_kmeansp2_siGL2n1_3RepConservedExons_5cluster_Liste_Didier_siPP_exons_down_1kbUpDown5pSS.png
#~ signals = sys.argv[ 2 ].split( ',' ) #~ [ 'siGL2', 'siPP' ]
typeEx_list = sys.argv[ 2 ].split( ',' )
out_dir = sys.argv[ 3 ]
makedirs( out_dir, exist_ok=True )
comp_cond = sys.argv[ 4 ].split( ',' )
metagene_dir = sys.argv[ 5 ] #~ ./summary_coverage_scaleHarm/metagene_mean
diff_metagene_dir = sys.argv[ 6 ] #~ summary_diffCoverage_scaleHarm/metagene_mean

prs = Presentation()

## some general variables
left_margin = Inches(0.1)
height_tbx = Inches(0.5)
width_tbx = Inches(0.5)
font_size = Pt(10)

space = 2.65
width = 2
width_sep = width*0.8

sign_htm = '_3RepConservedExons'

for typeEx in typeEx_list: #~ [ 'exons_down', 'exons_up', 'acceptors', 'donnors', 'AFE-promAl', 'ALE-polyAdAl', 'del-intron-ret', 'exons2' ]:
    for spSite in [ '3', '5' ]:
        ## slide with images ( barplot, heatmap over the three replicates and the three signals for clusters defined on the siGL2 signal )
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide( slide_layout )

        top_heatmap = Inches( 2.25 )
        top_metagene = top_heatmap

        # add a title
        slide.shapes.title.text = 'siPP ' + typeEx + ', 1kb +/- ' + spSite + '\'SS'

        # add the barplot of cluster conservation
        suffix_img = typeEx + '_1kbUpDown' + spSite + 'pSS.png'

        image_file = '/'.join( [ image_dir, 'barplot_siGL2_5clusterConservation_Liste_Didier_siPP_' + suffix_img ] )
        if not path.exists( image_file ):
            textbox = slide.shapes.add_textbox( left=Inches(3), top=Inches(3), width=Inches(1.5), height=height_tbx )
            textbox.text = 'No clustering available for this factor.'
            continue

        img = slide.shapes.add_picture( image_file, left=Inches(0.1), top=top_heatmap, width=Inches( width ) )

        left_heatmap = ( img.width/Inches(1) ) + 0.25
        #~_left_rownames = ( img.width/Inches(1) ) + 0.7
        left_rownames = left_heatmap - 0.55


        # add the heatmap for siGL2 coverage signal
        # add the heatmaps of the line
        for xxx in range(3):
            image_file = '/'.join( [ image_dir, 'heatmap_kmeansp2_' + comp_cond[ 1 ] + 'n' + str(xxx+1) + sign_htm + '_5cluster_Liste_Didier_siPP_' + suffix_img ] )
            img = slide.shapes.add_picture( image_file, left=Inches(left_heatmap + width_sep*( xxx )), top=top_heatmap, width=Inches( width ) )
            pass
        # add the rowname
        textbox = slide.shapes.add_textbox( left=Inches( left_rownames ), top=top_heatmap + (img.height * 0.3), width=Inches(1.5), height=height_tbx )
        textbox.text = comp_cond[ 1 ]
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        textbox.rotation = -90

        top_heatmap = top_heatmap + img.height

        # add the heatmap for siPP coverage signal
        # add the heatmaps of the line
        for xxx in range(3):
            image_file = '/'.join( [ image_dir, 'treatCov', 'heatmap_5clusteronsiGL2_' + comp_cond[ 0 ] + 'n' + str(xxx+1) + sign_htm + '_Liste_Didier_siPP_' + suffix_img ] )
            img = slide.shapes.add_picture( image_file, left=Inches(left_heatmap + width_sep*( xxx )), top=top_heatmap, width=Inches( width ) )
            pass
        # add the rowname
        textbox = slide.shapes.add_textbox( left=Inches( left_rownames ), top=top_heatmap + (img.height * 0.3), width=Inches(1.5), height=height_tbx )
        textbox.text = comp_cond[ 0 ]
        textbox.text_frame.word_wrap = True
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        textbox.rotation = -90

        top_heatmap = top_heatmap + img.height

        # add the heatmap for differential ( siPP-siGL2 ) coverage signal
        # add the heatmaps of the line
        for xxx in range(3):
            image_file = '/'.join( [ image_dir, 'diffCov', 'heatmap_5clusteronsiGL2_' + comp_cond[ 0 ] + 'n' + str(xxx+1) + '-' + comp_cond[ 1 ] + 'n' + str(xxx+1) + sign_htm + '_Liste_Didier_siPP_' + suffix_img ] )
            img = slide.shapes.add_picture( image_file, left=Inches(left_heatmap + width_sep*( xxx )), top=top_heatmap, width=Inches( width ) )
            pass
        # add the rowname
        textbox = slide.shapes.add_textbox( left=Inches( left_rownames ), top=top_heatmap + (img.height * 0.3), width=Inches(1.5), height=height_tbx )
        textbox.text = '-'.join( comp_cond )
        textbox.text_frame.word_wrap = True
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        textbox.rotation = -90

        top_heatmap = top_heatmap + img.height

        # add the names of the column for the table of heatmaps
        for xxx in range(3):
            textbox = slide.shapes.add_textbox( left=Inches(left_heatmap + width_sep*( xxx ) ), top=top_heatmap, width=Inches( 1.5 ), height=height_tbx )
            textbox.text = 'n%s' % xxx
            textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
            textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            textbox.left = int( textbox.left + Inches(width)/2 - textbox.width/2 )
            pass

        ## add the metagene plot of coverage signal
        left_metagene = img.left + img.width
            # coverage per condition
        textbox = slide.shapes.add_textbox( left=left_metagene, top=top_metagene , width=Inches(1.5), height=height_tbx )
        textbox.text = 'metagene mean cov. scaleHarm'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        image_file = '/'.join( [ metagene_dir, 'metagene_mean_MCF7_Liste_Didier_siPP_' + suffix_img ] )
        img = slide.shapes.add_picture( image_file, left=left_metagene, top=top_metagene + textbox.height, width=Inches( width ) )
        top_metagene = img.top + img.height + Inches(0.2)
            # differential coverage
        textbox = slide.shapes.add_textbox( left=left_metagene, top=top_metagene, width=Inches(1.5), height=height_tbx )
        textbox.text = 'metagene mean diff cov. scaleHarm'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        image_file = '/'.join( [ diff_metagene_dir, 'metagene_mean_diffCov_MCF7_Liste_Didier_siPP_' + suffix_img ] )
        img = slide.shapes.add_picture( image_file, left=left_metagene, top=top_metagene + textbox.height, width=Inches( width ) )
            
        pass
    pass


# save the presentation
out_file = out_dir + '/kmean_siGL2_on_otherSignals.pptx'
print( out_file )
prs.save( out_file )

