#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from math import floor

## public_signal page
inter_img_space_y = 0 #Inches( 0.1 )
offset_x = Inches( 0.25 )
public_signal_img_height = floor( 2 * img_height / 3 )
second_column_x = prs.slide_width - slide_marge_right - 2 * public_signal_img_height

column_dic = { 0:slide_marge_left, 1:second_column_x }

inEx = cov_in_exon #'50'
outOff = cov_out_off #'200'

cov_dir = fig_dir + '/coverage_summary/inEx%s_outOff%s' % ( inEx, outOff )
metagene_dir = '/cov/metagene_mean'
heatmap_dir = '/cov/5clusters/5mms'

# recover the list of directory for each factor
cov_paths = [ xxx for xxx in sorted( walk( cov_dir, followlinks=True ).__next__()[ 1 ] ) if xxx.startswith('public_') ]
#cov_path_ids = [ ( xxx.split( '_' )[ 1 ], xxx.split( '_' )[ 3 ] ) for xxx in cov_paths ]

# [ 'CTCF','H3K27ac','H3K27me3','H3K4me1','H3K4me2','H3K4me3','MBD3','Pol2' ]
for xxx, cov_path in enumerate( cov_paths ):
    print( cov_path )
    factor=cov_path.split( '_' )[ 1 ]
    sra=cov_path.split( 'ChIP-Seq_' )[ 1 ] #cov_path.split( 'mcf7_' )[ 2 ].split( '_treat_pileup' )[ 0 ]
    factor_dir = cov_dir + '/public_%s_ChIP-Seq_%s' % ( factor, sra )

    ## initialize position of the cursor
    colonne = column_dic[ xxx % 2 ]
    current_y = slide_marge_top
    current_x = colonne

    # initialize the page
    if xxx % 2 == 0:
        page_title = 'public_factor'
        file_path = script_dir + '/slide_creator.py'
        exec( exec_text )
        pass

    ## exon_scheme
    file_path = script_dir + '/splicing_site_region_scheme.py'
    exec( exec_text )

    fig_part_y = current_y + height_tbx + Inches( 0.3 )

    ######
    # simple base fraction
    #   # Metagene
    current_x = colonne
    current_y = fig_part_y + Inches( 0.2 )
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.66 ), height=height_tbx )
    textbox.text = factor + '_' + sra
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    textbox.text_frame.paragraphs[0].runs[0].font.bold = True

    #   # splicing sites titles
    current_y = current_y + height_tbx
    textbox = slide.shapes.add_textbox( left=current_x + offset_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '3\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    textbox = slide.shapes.add_textbox( left=current_x + offset_x + public_signal_img_height, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    textbox.text = '5\'SS'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 3'SS
    current_y = current_y + height_tbx
    image_file = factor_dir + metagene_dir + '/metagene_mean_' + exon_list_name + '_3SS.png'
    print( image_file )
    try:
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 5'SS
    current_x = current_x + img.width
    image_file = factor_dir + metagene_dir + '/metagene_mean_' + exon_list_name + '_5SS.png'
    try:
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small


    #   # Heatmap
    #   #   # up exons
    current_y = current_y + img.height + inter_img_space_y
    current_x = colonne
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
    textbox.text = 'Up exons'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 3'SS
    current_y = current_y + height_tbx
    try:
        image_file = factor_dir + heatmap_dir + '/heatmap_kmeansp2_publicn1_allExons_exon_up_' + exon_list_name + '_3SS.png'
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 5'SS
    current_x = current_x + public_signal_img_height
    try:
        image_file = factor_dir + heatmap_dir + '/heatmap_kmeansp2_publicn1_allExons_exon_up_' + exon_list_name + '_5SS.png'
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   # down exons
    current_x = colonne
    current_y = current_y + public_signal_img_height + inter_img_space_y
    textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
    textbox.text = 'Down exons'
    textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 3'SS
    current_y = current_y + height_tbx
    try:
        image_file = factor_dir + heatmap_dir + '/heatmap_kmeansp2_publicn1_allExons_exon_down_' + exon_list_name + '_3SS.png'
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small

    #   #   #   # 5'SS
    current_x = current_x + public_signal_img_height
    try:
        image_file = factor_dir + heatmap_dir + '/heatmap_kmeansp2_publicn1_allExons_exon_down_' + exon_list_name + '_5SS.png'
        img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    except:
        textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
        textbox.text = 'no image'
        textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small



    # #   # MBD-Seq metagene
    # current_y = fig_part_y + Inches( 0.2 )
    # current_x = second_column_x
    # textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.49 ), height=height_tbx )
    # textbox.text = 'MBD'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    # textbox.text_frame.paragraphs[0].runs[0].font.bold = True
    #
    # #   # splicing sites titles
    # current_y = current_y + height_tbx
    # textbox = slide.shapes.add_textbox( left=current_x + offset_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    # textbox.text = '3\'SS'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    #
    # textbox = slide.shapes.add_textbox( left=current_x + offset_x + public_signal_img_height, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    # textbox.text = '5\'SS'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    #
    # #   #   #   # 3'SS
    # current_y = current_y + height_tbx
    # image_file = MBD_dir + metagene_dir + '/metagene_mean_' + exon_list_name + '_3SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    # #   #   #   # 5'SS
    # current_x = current_x + img.width
    # image_file = MBD_dir + metagene_dir + '/metagene_mean_' + exon_list_name + '_5SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    #
    # #   # MBD-Seq heatmap
    # #   #   # up exons
    # current_y = current_y + img.height + inter_img_space_y
    # current_x = second_column_x
    # textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
    # textbox.text = 'Up exons'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    #
    # #   #   #   # 3'SS
    # current_y = current_y + height_tbx
    # image_file = MBD_dir + heatmap_dir + '/heatmap_kmeansp2_siGL2n2_allExons_exon_up_' + exon_list_name + '_3SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    # #   #   #   # 5'SS
    # current_x = current_x + img.width
    # image_file = MBD_dir + heatmap_dir + '/heatmap_kmeansp2_siGL2n2_allExons_exon_up_' + exon_list_name + '_5SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    # #   #   # down exons
    # current_x = second_column_x
    # current_y = current_y + img.height + inter_img_space_y
    # textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 1.55 ), height=height_tbx )
    # textbox.text = 'Down exons'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    #
    # #   #   #   # 3'SS
    # current_x = second_column_x
    # current_y = current_y + height_tbx
    # image_file = MBD_dir + heatmap_dir + '/heatmap_kmeansp2_siGL2n2_allExons_exon_down_' + exon_list_name + '_3SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    # #   #   #   # 5'SS
    # current_x = current_x + img.width
    # image_file = MBD_dir + heatmap_dir + '/heatmap_kmeansp2_siGL2n2_allExons_exon_down_' + exon_list_name + '_5SS.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    #
    # #   # RRBS metagene
    # current_y = current_y + img.height + Inches( 1 )
    # current_x = slide_marge_left
    # textbox = slide.shapes.add_textbox( left=current_x, top=current_y, width=Inches( 0.3 ), height=height_tbx )
    # textbox.text = 'RRBS'
    # textbox.text_frame.paragraphs[0].runs[0].font.size = font_size_small
    # textbox.text_frame.paragraphs[0].runs[0].font.bold = True
    #
    # #~ #   #   # up exons
    # #   #   #   # 3'SS
    # current_y = current_y + height_tbx
    # image_file = RRBS_dir + '/' + exon_list_name + '_3SS_]87:]_fract.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )
    #
    # #   #   #   # 5'SS
    # current_x = current_x + img.width
    # image_file = RRBS_dir + '/' + exon_list_name + '_5SS_]87:]_fract.png'
    # img = slide.shapes.add_picture( image_file, left=current_x, top=current_y, height=public_signal_img_height )

    pass
