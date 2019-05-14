#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

## splicing_site_region_scheme
init_left = Inches( 2.5 )
init_top = slide_marge_top
init_width = Inches( 1 )
init_height = Inches( 0.35 )
textbox_width = Inches( 1 )
inter_scheme_space = Inches( 1 )

shapes = slide.shapes
left = init_left
top = init_top
intron_width = init_width
exon_width = init_width
height = init_height


# 3'SS part
top = top + height / 2
height = 0
shape = shapes.add_shape(
    MSO_SHAPE.LINE_INVERSE, left, top, intron_width, height
)
shape.text = 'intron\n'
for xxx in range( len( shape.text_frame.paragraphs[0].runs ) ):
    shape.text_frame.paragraphs[0].runs[xxx].font.size = Pt( 14 )
    shape.text_frame.paragraphs[0].runs[xxx].font.color.rgb = RGBColor( 0, 0, 0 )
    pass

textbox = shapes.add_textbox( left - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '-' + outOff + 'bp'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )

left += intron_width
top = init_top
height = init_height
shape = shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, exon_width, height
)
shape.text = 'exon'
for xxx in range( len( shape.text_frame.paragraphs[0].runs ) ):
    shape.text_frame.paragraphs[0].runs[xxx].font.size = Pt( 14 )
    pass

textbox = shapes.add_textbox( left - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '3\'SS'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )

textbox = shapes.add_textbox( left + exon_width - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '+' + inEx + 'bp'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )



# 5'SS part
left += exon_width + inter_scheme_space
top = init_top
height = init_height
shape = shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, exon_width, height
)
shape.text = 'exon'
for xxx in range( len( shape.text_frame.paragraphs[0].runs ) ):
    shape.text_frame.paragraphs[0].runs[xxx].font.size = Pt( 14 )
    pass

textbox = shapes.add_textbox( left - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '-' + inEx + 'bp'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )

textbox = shapes.add_textbox( left + exon_width - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '5\'SS'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )

left += exon_width
top = top + height / 2
height = 0
shape = shapes.add_shape(
    MSO_SHAPE.LINE_INVERSE, left, top, intron_width, height
)
shape.text = 'intron\n'
for xxx in range( len( shape.text_frame.paragraphs[0].runs ) ):
    shape.text_frame.paragraphs[0].runs[xxx].font.size = Pt( 14 )
    shape.text_frame.paragraphs[0].runs[xxx].font.color.rgb = RGBColor( 0, 0, 0 )
    pass

textbox = shapes.add_textbox( left + intron_width - textbox_width / 2, top + height, width = textbox_width, height = Inches( 0.3 ) )
textbox.text = '+' + outOff + 'bp'
textbox.text_frame.paragraphs[0].runs[0].font.size = Pt( 14 )
