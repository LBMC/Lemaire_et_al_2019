#!/usr/bin/python3

from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import sys
from os import makedirs,path,walk,environ
from subprocess import getoutput
from math import floor
import glob

## Usage: [ python3 ] pptx_features_wNbs_builder.py <suffix> <res_dir> [ --out PATH ]

## define a text of commands (pseudo-functionalize execution of a file)
exec_text = '''with open( file_path, 'r' ) as fich:
    fich_text = ''.join( [ xxx for xxx in fich ] )
    pass
exec( fich_text )
'''

def legend_hide( img, slide ):
    shapes = slide.shapes
    shape = shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left=img.left, top=img.top + img.height - Inches( 0.5 ), width=img.width, height=Inches( 0.5 )
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor( 255, 255, 255 )
    shape.line.fill.background()

def legend_frame():
    colors = [ ( 102, 102 ,102 ), ( 153, 153, 153 ), ( 204, 204, 204 ), ( 255, 178, 102 ), ( 255, 128, 0 ), ( 204, 102, 0 ), ( 255, 102, 178 ), ( 204, 0, 102 ), ( 255, 0, 0 ), ( 0, 136, 255 ) ]
    list_names = [ 'psi_100_95_FDB', 'psi_60_40_FDB', 'psi_5_0_FDB', 'intron_50-200', 'intron_200-500', 'intron_1000-4000', 'GC_25-35', 'GC_55-65', 'exons up', 'exons down' ]

    leg_y = Inches( 5.5 )
    line_inter = Inches( 0.2 )

    for xxx in range( len( colors ) ):
        line_y = leg_y + line_inter + xxx*( line_inter )
        shapes = slide.shapes
        shape = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=slide_marge_left, top=line_y, width=Inches( 0.2 ), height=Inches( 0.1 )
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor( colors[ xxx ][ 0 ], colors[ xxx ][ 1 ], colors[ xxx ][ 2 ] )
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = RGBColor( colors[ xxx ][ 0 ], colors[ xxx ][ 1 ], colors[ xxx ][ 2 ] )

        txt_height = Inches( 0.2 )
        txt_width = Inches( 1.5 )
        textbox = slide.shapes.add_textbox( left=slide_marge_left + Inches( 0.1 ), top=line_y - ( txt_height/4 ), width=txt_width, height=txt_height )
        p = textbox.text_frame.paragraphs[0]
        p.text = list_names[ xxx ]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt( 8 )
        pass

def get_path( glob_pattern ):
    try:
        right_path = glob.glob( glob_pattern )[ 0 ]
    except:
        print( 'Error! Cannot find: ' + glob_pattern, file=sys.stderr )
        sys.exit()

    return( right_path )

#######


## define directory variables
# Usage: python3 ../src/pdf/pptx_covareg_multiAnal_builder.py <out_suffix> <res_dir1>,<res_dir2>,... <inEx>,<outOff>,<window> <measure1>,<measure2>,... [ --out <path> ] [ --res-name <name1>,<name2>,... ]
# out_suffix: string added at the end of the pptx basename
# res_dirX: directory of Felrs results (ex: siDNMT3b/results )
# measureX: measure to display ( GpGpG, CpG, T, ... )
# path: directory in which to write the produced pptx [ default is '.' ]
#
## Example
# wdir='/home/sebastien/analyses_SEBASTIEN/analyse_RNA-Seq_methChromRegulators-siPP-siGL2/Felrs_DNMT3b_advanced/query_exon_list'
# python3 \
#   pptx_coverage_multiAnal_builder.py \
#   ed3B_Nbs \
#   ${wdir}/siDNMT3b-siGL2_nbup,${wdir}/siDNMT3b-siGL2,${wdir}/siDNMT3b-siGL2_nbdown \
#   50,200 \
#   MBD_Seq,MNase_Seq \
#   coverage \
#   --out ${wdir}/../pptx_reports \
#   --res-name n-1,n,n+1 \
#   ;

suf = sys.argv[ 1 ]
res_dir_list = sys.argv[ 2 ].split( ',' )
for res_dir in res_dir_list:
    if res_dir.startswith( '~/' ):
        res_dir = environ[ 'HOME' ] + res_dir[ 1: ]
        pass
    pass

scan_region = sys.argv[ 3 ].split( ',' )
out_dir='.'
if '--out' in sys.argv:
    out_dir = sys.argv[ sys.argv.index( '--out' ) + 1 ]
    pass
script_dir = path.dirname( path.realpath( __file__ ) )
sub_fig_dir = 'results/figures'

exon_pos_list = [ str( xxx ) for xxx in range( len( res_dir_list ) ) ]
if '--res-name' in sys.argv:
    exon_pos_list = sys.argv[ sys.argv.index( '--res-name' ) + 1 ].split( ',' )
    pass

center_meas = False
if '--center' in sys.argv:
    center_meas = True
    pass

# measure_list = [ 'exon_length', 'intron_length', 'force_acceptor', 'force_donor' ]
measure_list = sys.argv[ 4 ].split( ',' )
measure_branch = sys.argv[ 5 ] # 'coverage, base_fraction, mCpG, exon_feature, ...'
print( 'Type of measure: ' + measure_branch, file=sys.stderr )

measure_type_dict = {
'exon_feature': [''],
'coverage': [ 'median', 'mean' ],
'coverage_heatmap': [''],
'coverage_max_repart': [''],
'base_fraction': [''],
'base_fraction_heatmap': [''],
'base_fraction_repart': [''],
'weblogo': [''],
'methCpG': [''],
'methCpG_meta': [''],
'rna_folding': ['']
}

page_title_dict = {
'exon_feature': 'Exon Feature',
'coverage': 'Coverage',
'coverage_heatmap': 'Coverage Heatmap',
'coverage_max_repart': 'Coverage Max Repart',
'base_fraction': 'Base Fraction',
'base_fraction_heatmap': 'Base Fraction Heatmap',
'base_fraction_repart': 'Base Fraction Repart',
'methCpG': '',
'methCpG_meta': 'methCpG_metaplot',
'weblogo': 'WebLogo',
'rna_folding': 'RNA Folding Free Energy'
}

for measure_type in measure_type_dict[ measure_branch ]:
    pptx_name = out_dir + '/pptx_' + measure_branch + '_multiAnal_'
    if measure_branch == 'coverage':
        pptx_name = pptx_name + measure_type + '_'
        pass
    pptx_name = pptx_name + suf

    ## initialize pptx
    prs = Presentation()
    prs.slide_width = Inches( 11.69 )
    prs.slide_height = Inches( 10.27 )
    # prs.slide_width = Inches( 11.69 )
    # prs.slide_height = Inches( 8.27 )

    slide_marge_left = Inches( 0.39 )
    slide_marge_right = slide_marge_left
    slide_marge_top = slide_marge_left + Inches( 0.1 )
    slide_marge_bottom = slide_marge_left

    height_tbx = Inches( 0.3 )
    img_height = Inches( 3.5 ) #2.45
    font_size_big = Pt( 24 ) #16
    font_size_small = Pt( 12 )

    ## pages
    #########
    # exon features for exon of interest and neighbour exons
    feature_dir_dict = {
    }

    if measure_branch == 'exon_feature':
        feature_dir = sub_fig_dir + '/exon_feature'
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/exon_feature_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'base_fraction':
        feature_dir = sub_fig_dir + '/base_fraction/inEx%s_outOff%s/win%sb' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/base_fraction_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'base_fraction_heatmap':
        feature_dir = sub_fig_dir + '/base_fraction/inEx%s_outOff%s/win%sb/heatmaps' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/base_fraction_heatmap_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'base_fraction_repart':
        feature_dir = sub_fig_dir + '/base_fraction/inEx%s_outOff%s/win%sb/reparts' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/base_fraction_repart_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'weblogo':
        feature_dir = sub_fig_dir + '/base_fraction/inEx%s_outOff%s/win%sb/weblogos' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        exon_list_dic = {}
        for res_dir in res_dir_list:
            exon_list_dic[ res_dir ] = sorted( [ xxx.split( '/' )[ -1 ].split( '_3SS_' )[ 0 ] for xxx in glob.glob( res_dir + '/' + feature_dir + '/*_3SS_*.png' ) ] )
            pass
        exon_list_nb = len( exon_list_dic[ res_dir ] )
        print( )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/weblogo_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'coverage':
        if center_meas:
            feature_dir = sub_fig_dir + '/coverage_summary/aroundCenter%s' % ( scan_region[ 0 ] )
        else:
            feature_dir = sub_fig_dir + '/coverage_summary/inEx%s_outOff%s' % ( scan_region[ 0 ], scan_region[ 1 ] )
            pass
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/coverage_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'coverage_heatmap':
        feature_dir = sub_fig_dir + '/coverage_summary/inEx%s_outOff%s' % ( scan_region[ 0 ], scan_region[ 1 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/coverage_heatmap_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'coverage_max_repart':
        feature_dir = sub_fig_dir + '/coverage_summary/inEx%s_outOff%s' % ( scan_region[ 0 ], scan_region[ 1 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/coverage_max_repart_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'methCpG':
        feature_dir = sub_fig_dir + '/methCpG_meta/inEx%s_outOff%s/win%sb' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        for stat in [ 'nbCpG', 'mCpG_dist', 'CpG-cat' ]:
            page_title = 'CpG Methylation'
            file_path = script_dir + '/methCpG_multiAnal.py'
            exec( exec_text )
            pass
        pass
    elif measure_branch == 'methCpG_meta':
        feature_dir = sub_fig_dir + '/methCpG_meta/inEx%s_outOff%s/win%sb' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        page_title = 'CpG Methylation Metaplot'
        file_path = script_dir + '/methCpG_meta_multiAnal.py'
        exec( exec_text )
        pass
    elif measure_branch == 'rna_folding':
        feature_dir = sub_fig_dir + '/rna_folding/inEx%s_outOff%s/win%sb' % ( scan_region[ 0 ], scan_region[ 1 ], scan_region[ 2 ] )
        page_title = page_title_dict[ measure_branch ]
        file_path = script_dir + '/rna_folding_multiAnal.py'
        exec( exec_text )
        pass
    else:
        print( 'Bad type of measure. No PPTX produced', file=sys.stderr )
        pass


    #########


    out_file = pptx_name + '.pptx'
    makedirs( path.dirname( out_file ), exist_ok=True )
    print( out_file )
    prs.save( out_file )

    pass
